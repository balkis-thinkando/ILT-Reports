
#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
ppt_audit.py
-------------
Scan a directory for PowerPoint files (.pptx/.pptm) and collect useful metadata:
- slide_count
- picture_count
- estimated_screenshots (heuristic based on image dimensions/aspect ratio)
- text statistics (word count, char count)
- keywords (naive frequency-based)
- short_summary (first few meaningful lines across slides)
- course_name (big title on the first slide, **appended with platform** when detected)
- platform (Cloud/DC when detected from file name)

Outputs a CSV (and optional JSON).
"""

import argparse
import io
import json
import sys
import re
from collections import Counter
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER
from PIL import Image
import pandas as pd


STOPWORDS = {
    # minimal stopword set; feel free to extend
    "the","and","for","that","with","this","from","your","have","will","into","are","was","were","has","had","but",
    "not","you","our","out","any","can","all","its","it's","they","them","then","than","these","those","over","about",
    "what","when","where","which","while","within","without","between","into","onto","off","too","also","more","less",
    "how","why","who","whom","whose","a","an","in","on","at","of","by","to","as","or","is","be","we","it"
}

def condense(s: str, max_chars: int = 200) -> str:
    s = " ".join((s or "").split())
    if len(s) > max_chars:
        s = s[: max_chars - 1] + "…"
    return s

def is_probable_screenshot(width, height):
    """
    Heuristic: treat large, screen-like images as 'screenshots'.
    - At least 1000 px wide and 600 px tall
    - Aspect ratio close to common screens (16:9 ~= 1.78, 3:2 ~= 1.5, 4:3 ~= 1.33)
    """
    if width is None or height is None:
        return False
    if width < 1000 or height < 600:
        return False
    ar = width / float(height)
    # within +/- 0.2 of common aspect ratios
    for target in (1.78, 1.6, 1.5, 1.33):
        if abs(ar - target) <= 0.2:
            return True
    return False


def extract_text_from_slide(slide):
    """Gather visible text from text frames and tables."""
    lines = []
    for shape in slide.shapes:
        # text frames
        if hasattr(shape, "has_text_frame") and shape.has_text_frame:
            txt = shape.text or ""
            if txt:
                for line in txt.splitlines():
                    line = line.strip()
                    if line:
                        lines.append(line)
        # tables
        if hasattr(shape, "has_table") and shape.has_table:
            tbl = shape.table
            for r in tbl.rows:
                for c in r.cells:
                    cell_txt = (c.text or "").strip()
                    if cell_txt:
                        lines.append(cell_txt)
    return lines


def tokenize(text):
    words = re.findall(r"[A-Za-z][A-Za-z0-9'-]*", text)
    return [w.lower() for w in words]


def keywords_from_text(text, topn=12):
    words = [w for w in tokenize(text) if w not in STOPWORDS and len(w) >= 4]
    if not words:
        return []
    counts = Counter(words)
    return [w for w, _ in counts.most_common(topn)]


def short_summary_from_lines(lines, max_items=5, max_chars=140):
    summary = []
    seen = set()
    for line in lines:
        line_norm = " ".join(line.split())
        if not line_norm or line_norm.lower() in seen:
            continue
        seen.add(line_norm.lower())
        if len(line_norm) > max_chars:
            line_norm = line_norm[: max_chars - 1] + "…"
        summary.append(line_norm)
        if len(summary) >= max_items:
            break
    return " • ".join(summary)


def extract_course_name(prs: Presentation):
    """Return the big title from the first slide, if we can guess it."""
    try:
        slide = prs.slides[0]
    except IndexError:
        return None

    # 1) Built-in title placeholder (most reliable)
    try:
        title_shape = slide.shapes.title
    except Exception:
        title_shape = None
    if title_shape is not None:
        t = condense(getattr(title_shape, "text", "") or "")
        if t:
            return t

    # 2) Any placeholder explicitly marked as TITLE/CENTER_TITLE
    for shape in slide.shapes:
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                ph = shape.placeholder_format
                if getattr(ph, "type", None) in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                    t = condense(getattr(shape, "text", "") or "")
                    if t:
                        return t
        except Exception:
            pass

    # 3) Heuristic: longest text line in the top half of the slide
    H = getattr(prs, "slide_height", None) or 0
    best_line = ""
    best_score = -1
    for shape in slide.shapes:
        if hasattr(shape, "has_text_frame") and shape.has_text_frame:
            y_center = getattr(shape, "top", 0) + getattr(shape, "height", 0) / 2
            top_weight = 20 if (H and y_center < H / 2) else 0
            text = getattr(shape, "text", "") or ""
            for raw in (ln.strip() for ln in text.splitlines() if ln.strip()):
                line = condense(raw, 140)
                score = len(line) + top_weight
                if score > best_score:
                    best_score = score
                    best_line = line

    return best_line or None


def detect_platform_from_filename(file_name: str):
    """Return ('Cloud' or 'DC') if the file name contains -CLD- or -DC- (case-insensitive)."""
    up = file_name.upper()
    if "-CLD-" in up:
        return "Cloud"
    if "-DC-" in up:
        return "DC"
    return None


def analyze_ppt(path: Path):
    try:
        prs = Presentation(str(path))
    except Exception as e:
        return {"file_path": str(path), "error": f"open_failed: {e}"}

    slide_count = len(prs.slides)
    picture_count = 0
    screenshot_estimate = 0
    all_lines = []
    course_title = extract_course_name(prs) or Path(path).stem
    platform = detect_platform_from_filename(path.name)

    for slide in prs.slides:
        all_lines.extend(extract_text_from_slide(slide))

        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                picture_count += 1
                try:
                    img = shape.image
                    b = io.BytesIO(img.blob)
                    with Image.open(b) as im:
                        w, h = im.size
                    if is_probable_screenshot(w, h):
                        screenshot_estimate += 1
                except Exception:
                    # if we cannot inspect image, ignore
                    pass

    full_text = "\n".join(all_lines)
    words = tokenize(full_text)
    total_words = len(words)
    total_chars = len(full_text)

    keywords = keywords_from_text(full_text, topn=12)
    short_summary = short_summary_from_lines(all_lines, max_items=5, max_chars=140)

    # Compose course_name with platform when available
    if platform:
        course_name = f"{course_title} - {platform}"
    else:
        course_name = course_title  # leave unchanged if platform not detected

    stat = {
        "Course Name": course_name,
        "Platform": platform or "",
        "File Name": path.name,
        "Number of slides": slide_count,
        "Number of pictures": picture_count,
        "Estimated screenshots": screenshot_estimate,
        "Total words": total_words,
        "keywords": ", ".join(keywords),
        "Short Summary": short_summary
        
    }
    return stat



def find_ppts(
    root: Path,
    recursive: bool = True,
    include_pptm: bool = True,
    debug: bool = False,
):
    """
    Discover PowerPoint files under `root` (Dropbox-safe).
    - Uses os.walk to follow symlinks/junctions (common in Dropbox/team spaces)
    - Case-insensitive match for .pptx/.pptm/.ppt
    - Detects cloud placeholders and logs a count so you know what to hydrate
    - Skips Office temp/lock files (~$*.pptx)
    """
    import os, sys, stat

    exts = {".pptx", ".ppt"}  # include .ppt as well
    if include_pptm:
        exts.add(".pptm")

    root = root.expanduser()

    def looks_like_placeholder(p: Path) -> bool:
        # Heuristics for macOS File Provider / Windows CFAPI placeholders
        try:
            if p.is_file():
                try:
                    st = p.stat()
                    # Some placeholders report size 0
                    if st.st_size == 0:
                        return True
                except Exception:
                    pass
                return False
            # If not a dir but exists, could be a cloud placeholder
            if p.exists() and not p.is_dir():
                # Windows: check file attributes if exposed
                try:
                    st = p.stat(follow_symlinks=False)
                    attrs = getattr(st, "st_file_attributes", 0)
                    FILE_ATTRIBUTE_REPARSE = 0x0400
                    FILE_ATTRIBUTE_RECALL_ON_DATA_ACCESS = 0x00400000
                    if attrs & (FILE_ATTRIBUTE_REPARSE | FILE_ATTRIBUTE_RECALL_ON_DATA_ACCESS):
                        return True
                except Exception:
                    pass
                # Try opening; failure suggests placeholder
                try:
                    with open(p, "rb"):
                        pass
                    return False
                except Exception:
                    return True
        except Exception:
            return False
        return False

    real_files = []
    placeholders = []

    if recursive:
        walker = os.walk(root, topdown=True, followlinks=True)
    else:
        # emulate non-recursive walk
        walker = [(str(root), [d.name for d in root.iterdir() if d.is_dir()], [f.name for f in root.iterdir() if f.is_file() or f.exists()])]

    for dirpath, dirnames, filenames in walker:
        for fname in filenames:
            p = Path(dirpath) / fname
            # skip lock files
            if fname.startswith("~$"):
                if debug:
                    print(f"[skip temp] {p}", file=sys.stderr)
                continue
            # extension check
            try:
                suff = p.suffix.lower()
            except Exception:
                suff = ""
            if suff not in exts:
                continue

            # placeholder vs real
            try:
                if p.is_file():
                    real_files.append(p.resolve())
                else:
                    if looks_like_placeholder(p):
                        placeholders.append(p)
                    else:
                        # try a best-effort stat/open to hydrate on macOS; ignore errors
                        try:
                            _ = p.stat()
                            with open(p, "rb"):
                                pass
                            # If open works, treat as real
                            real_files.append(p.resolve())
                        except Exception:
                            placeholders.append(p)
            except OSError as e:
                if debug:
                    print(f"[skip error] {p} -> {e}", file=sys.stderr)
                continue

    uniq = sorted(set(real_files))
    if debug:
        import sys
        print(f"[debug] found {len(uniq)} 'real' PPT files under {root}", file=sys.stderr)
        if placeholders:
            print(f"[debug] detected {len(placeholders)} placeholders (online-only). Hydrate these folders or files.", file=sys.stderr)
            for pp in placeholders[:20]:
                print(f"[debug]   placeholder: {pp}", file=sys.stderr)
            if len(placeholders) > 20:
                print(f"[debug]   ... (+{len(placeholders)-20} more)", file=sys.stderr)

    return uniq


def main():
    p = argparse.ArgumentParser(description="Collect metadata & summaries from PowerPoint files.")
    p.add_argument("--root", required=True, help="Directory to scan")
    p.add_argument("--output", required=True, help="CSV output path (e.g., report.csv)")
    p.add_argument("--json", help="Optional JSON output path (e.g., report.json)")
    p.add_argument("--nonrecursive", action="store_true", help="Only scan the top-level directory")
    p.add_argument("--since-days", type=int, default=None, help="Only include files modified in the last N days")
    args = p.parse_args()

    root = Path(args.root).expanduser()
    if not root.exists() or not root.is_dir():
        print(f"[error] root directory does not exist: {root}", file=sys.stderr)
        sys.exit(2)

    recursive = not args.nonrecursive
    files = find_ppts(root, recursive=recursive, include_pptm=True)

    if args.since_days is not None:
        import time
        cutoff = time.time() - (args.since_days * 86400)
        files = [f for f in files if f.stat().st_mtime >= cutoff]

    rows = []
    for f in files:
        res = analyze_ppt(f)
        rows.append(res)

    df = pd.DataFrame(rows)
    df.to_csv(args.output, index=False)

    if args.json:
        # convert DataFrame rows to a list of dicts
        with open(args.json, "w", encoding="utf-8") as jf:
            json.dump(rows, jf, ensure_ascii=False, indent=2)

    print(f"Analyzed {len(files)} files.")
    print(f"CSV saved to: {args.output}")
    if args.json:
        print(f"JSON saved to: {args.json}")


if __name__ == "__main__":
    main()
